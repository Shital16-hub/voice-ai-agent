"""
Document processing module for loading, parsing, and chunking documents.
"""
import os
import re
import logging
import tempfile
from typing import List, Dict, Any, Optional, Tuple, Iterator, Union
from pathlib import Path
import hashlib

from .config import get_document_processor_config

logger = logging.getLogger(__name__)

class Document:
    """
    Document class representing a chunk of text with metadata.
    """
    
    def __init__(
        self,
        text: str,
        metadata: Optional[Dict[str, Any]] = None,
        doc_id: Optional[str] = None
    ):
        """
        Initialize Document.
        
        Args:
            text: The document text content
            metadata: Optional metadata about the document
            doc_id: Optional document ID
        """
        self.text = text
        self.metadata = metadata or {}
        self.doc_id = doc_id or self._generate_id()
    
    def _generate_id(self) -> str:
        """Generate a unique ID based on content."""
        content_hash = hashlib.md5(self.text.encode('utf-8')).hexdigest()
        return f"doc_{content_hash}"
    
    def __str__(self) -> str:
        return f"Document(id={self.doc_id}, text={self.text[:50]}..., metadata={self.metadata})"
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary."""
        return {
            "id": self.doc_id,
            "text": self.text,
            "metadata": self.metadata
        }
    
    @classmethod
    def from_dict(cls, data: Dict[str, Any]) -> 'Document':
        """Create from dictionary."""
        return cls(
            text=data["text"],
            metadata=data.get("metadata", {}),
            doc_id=data.get("id")
        )

class DocumentProcessor:
    """
    Process documents for knowledge base ingestion.
    """
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """
        Initialize DocumentProcessor.
        
        Args:
            config: Optional configuration dictionary
        """
        self.config = config or get_document_processor_config()
        self.chunk_size = self.config["chunk_size"]
        self.chunk_overlap = self.config["chunk_overlap"]
        self.max_document_size = self.config["max_document_size_mb"] * 1024 * 1024
        self.supported_types = self.config["supported_types"]
        
        logger.info(f"Initialized DocumentProcessor with chunk_size={self.chunk_size}, "
                  f"chunk_overlap={self.chunk_overlap}")
    
    def is_supported_file(self, file_path: str) -> bool:
        """
        Check if file type is supported.
        
        Args:
            file_path: Path to the file
            
        Returns:
            True if file type is supported
        """
        ext = os.path.splitext(file_path)[1].lower()
        return ext in self.supported_types
    
    def load_document(self, file_path: str) -> List[Document]:
        """
        Load and process a document file.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            List of Document objects
        """
        # Check if file exists
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        # Check if file type is supported
        if not self.is_supported_file(file_path):
            raise ValueError(f"Unsupported file type: {file_path}")
        
        # Check file size
        file_size = os.path.getsize(file_path)
        if file_size > self.max_document_size:
            raise ValueError(f"File too large: {file_size} bytes "
                          f"(max: {self.max_document_size} bytes)")
        
        # Extract file metadata
        file_metadata = self._extract_file_metadata(file_path)
        
        # Parse document based on file type
        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if ext in ['.txt', '.md']:
                text = self._load_text_file(file_path)
            elif ext == '.pdf':
                text = self._load_pdf(file_path)
            elif ext in ['.docx', '.doc']:
                text = self._load_word_document(file_path)
            elif ext in ['.pptx', '.ppt']:
                text = self._load_presentation(file_path)
            elif ext in ['.xlsx', '.xls', '.csv']:
                text = self._load_spreadsheet(file_path)
            elif ext in ['.html', '.htm']:
                text = self._load_html(file_path)
            else:
                raise ValueError(f"Unsupported file type: {ext}")
            
            # Chunk the document
            chunks = self._chunk_text(text)
            
            # Create Document objects
            documents = []
            for i, chunk in enumerate(chunks):
                # Create metadata for this chunk
                chunk_metadata = file_metadata.copy()
                chunk_metadata["chunk_index"] = i
                chunk_metadata["chunk_count"] = len(chunks)
                
                # Create Document
                doc = Document(
                    text=chunk,
                    metadata=chunk_metadata,
                    doc_id=f"{os.path.basename(file_path)}_{i}"
                )
                documents.append(doc)
            
            logger.info(f"Processed document {file_path}: created {len(documents)} chunks")
            return documents
            
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {e}")
            raise
    
    def load_text(self, text: str, source_name: str = "text_input") -> List[Document]:
        """
        Load and process text directly.
        
        Args:
            text: Text content
            source_name: Name to use as source identifier
            
        Returns:
            List of Document objects
        """
        # Create basic metadata
        metadata = {
            "source": source_name,
            "source_type": "direct_text",
            "file_path": None,
            "file_type": None,
            "file_name": None,
            "created_at": None,
            "modified_at": None
        }
        
        # Chunk the text
        chunks = self._chunk_text(text)
        
        # Create Document objects
        documents = []
        for i, chunk in enumerate(chunks):
            # Create metadata for this chunk
            chunk_metadata = metadata.copy()
            chunk_metadata["chunk_index"] = i
            chunk_metadata["chunk_count"] = len(chunks)
            
            # Create Document
            doc = Document(
                text=chunk,
                metadata=chunk_metadata,
                doc_id=f"{source_name}_{i}"
            )
            documents.append(doc)
        
        logger.info(f"Processed text input '{source_name}': created {len(documents)} chunks")
        return documents
    
    def _extract_file_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        Extract metadata from file.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Dictionary with metadata
        """
        file_stat = os.stat(file_path)
        return {
            "source": os.path.basename(file_path),
            "source_type": "file",
            "file_path": os.path.abspath(file_path),
            "file_type": os.path.splitext(file_path)[1].lower(),
            "file_name": os.path.basename(file_path),
            "file_size": file_stat.st_size,
            "created_at": file_stat.st_ctime,
            "modified_at": file_stat.st_mtime
        }
    
    def _chunk_text(self, text: str) -> List[str]:
        """
        Split text into chunks.
        
        Args:
            text: Text to chunk
            
        Returns:
            List of text chunks
        """
        # Clean text
        text = self._clean_text(text)
        
        # Simple chunking by characters with overlap
        chunks = []
        start = 0
        
        while start < len(text):
            # Get chunk
            end = min(start + self.chunk_size, len(text))
            
            # If not the last chunk, try to find a good break point
            if end < len(text):
                # Look for paragraph break first
                paragraph_break = text.rfind('\n\n', start, end)
                if paragraph_break != -1 and paragraph_break > start + self.chunk_size // 2:
                    end = paragraph_break + 2
                else:
                    # Look for line break
                    line_break = text.rfind('\n', start, end)
                    if line_break != -1 and line_break > start + self.chunk_size // 2:
                        end = line_break + 1
                    else:
                        # Look for sentence break
                        sentence_break = max(
                            text.rfind('. ', start, end),
                            text.rfind('? ', start, end),
                            text.rfind('! ', start, end)
                        )
                        if sentence_break != -1 and sentence_break > start + self.chunk_size // 2:
                            end = sentence_break + 2
                        else:
                            # Look for word break
                            word_break = text.rfind(' ', start, end)
                            if word_break != -1 and word_break > start + self.chunk_size // 2:
                                end = word_break + 1
            
            # Add chunk
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position for next chunk
            start = end - self.chunk_overlap
            
            # Ensure progress
            if start < end - self.chunk_overlap:
                start = end
        
        return chunks
    
    def _clean_text(self, text: str) -> str:
        """
        Clean text for processing.
        
        Args:
            text: Text to clean
            
        Returns:
            Cleaned text
        """
        # Replace multiple spaces with single space
        text = re.sub(r'\s+', ' ', text)
        
        # Replace multiple newlines with double newline
        text = re.sub(r'\n+', '\n\n', text)
        
        # Remove excessive whitespace at start and end
        text = text.strip()
        
        return text
    
    def _load_text_file(self, file_path: str) -> str:
        """Load text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # Try different encodings
            encodings = ['latin-1', 'windows-1252', 'iso-8859-1']
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            
            # If all failed, try binary and replace errors
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()
    
    def _load_pdf(self, file_path: str) -> str:
        """Load PDF file."""
        try:
            import fitz  # PyMuPDF
            
            text = ""
            with fitz.open(file_path) as doc:
                for page in doc:
                    text += page.get_text()
            
            return text
            
        except ImportError:
            logger.warning("PyMuPDF not installed. Falling back to pdfplumber.")
            try:
                import pdfplumber
                
                text = ""
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        text += page.extract_text() or ""
                
                return text
                
            except ImportError:
                raise ImportError("Neither PyMuPDF nor pdfplumber is installed. "
                               "Please install one of them for PDF support.")
    
    def _load_word_document(self, file_path: str) -> str:
        """Load Word document."""
        try:
            import docx
            
            doc = docx.Document(file_path)
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            
            return '\n'.join(full_text)
            
        except ImportError:
            try:
                # Try using mammoth if python-docx isn't available
                import mammoth
                
                with open(file_path, 'rb') as f:
                    result = mammoth.extract_raw_text(f)
                
                return result.value
                
            except ImportError:
                raise ImportError("Neither python-docx nor mammoth is installed. "
                               "Please install one of them for DOCX support.")
    
    def _load_presentation(self, file_path: str) -> str:
        """Load PowerPoint presentation."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            
            return '\n\n'.join(text)
            
        except ImportError:
            raise ImportError("python-pptx is not installed. "
                           "Please install it for PowerPoint support.")
    
    def _load_spreadsheet(self, file_path: str) -> str:
        """Load Excel or CSV file."""
        ext = os.path.splitext(file_path)[1].lower()
        
        if ext == '.csv':
            try:
                import pandas as pd
                
                df = pd.read_csv(file_path)
                return df.to_string(index=False)
                
            except ImportError:
                # Fallback to csv module
                import csv
                
                rows = []
                with open(file_path, 'r', newline='', encoding='utf-8') as csvfile:
                    reader = csv.reader(csvfile)
                    for row in reader:
                        rows.append(','.join(row))
                
                return '\n'.join(rows)
                
        else:  # Excel files
            try:
                import pandas as pd
                
                # Read all sheets
                xlsx = pd.ExcelFile(file_path)
                texts = []
                
                for sheet_name in xlsx.sheet_names:
                    df = pd.read_excel(xlsx, sheet_name=sheet_name)
                    texts.append(f"Sheet: {sheet_name}")
                    texts.append(df.to_string(index=False))
                
                return '\n\n'.join(texts)
                
            except ImportError:
                raise ImportError("pandas is not installed. "
                               "Please install it for Excel support.")
    
    def _load_html(self, file_path: str) -> str:
        """Load HTML file."""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f, 'html.parser')
                
                # Remove script and style elements
                for script in soup(["script", "style"]):
                    script.extract()
                
                # Get text
                text = soup.get_text()
                
                # Break into lines and remove leading/trailing whitespace
                lines = (line.strip() for line in text.splitlines())
                
                # Break multi-headlines into a line each
                chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                
                # Remove blank lines
                text = '\n'.join(chunk for chunk in chunks if chunk)
                
                return text
                
        except ImportError:
            raise ImportError("beautifulsoup4 is not installed. "
                           "Please install it for HTML support.")